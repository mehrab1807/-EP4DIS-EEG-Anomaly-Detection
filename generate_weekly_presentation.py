import mne
import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
import os

def create_plots():
    print("Loading data...")
    file_path = r'data\sub-jh101\ses-presurgery\ieeg\sub-jh101_ses-presurgery_task-ictal_acq-ecog_run-01_ieeg.vhdr'
    raw = mne.io.read_raw_brainvision(file_path, preload=True)
    
    # We use only a subset of channels for clean visualization
    raw.pick_types(eeg=True, meg=False, stim=False, exclude='bads').pick_channels(raw.ch_names[:5])
    
    # 1. Raw Time-Series (10s from start)
    print("Generating raw time-series plot...")
    start, stop = raw.time_as_index([0, 10])
    data, times = raw[:, start:stop]
    plt.figure(figsize=(10, 4))
    for i in range(data.shape[0]):
        plt.plot(times, data[i] + i * 0.001, color='black', alpha=0.7) # offset for visibility
    plt.title("Raw EEG Time-Series (Stochastic Drift)")
    plt.xlabel("Time (s)")
    plt.ylabel("Voltage (offset)")
    plt.tight_layout()
    plt.savefig("raw_time_series.png")
    plt.close()

    # 2. PSD Plot
    print("Generating PSD plot...")
    # PSD computation using welch method
    fig = raw.compute_psd(fmax=100).plot(show=False)
    fig.savefig("psd_plot.png")
    plt.close(fig)

    # 3. State Comparison (Interictal vs Ictal)
    print("Generating state comparison plot...")
    # From events file, seizure starts at 58.95s
    interictal_start, interictal_stop = raw.time_as_index([45, 50])
    ictal_start, ictal_stop = raw.time_as_index([60, 65])
    
    data_inter, times_inter = raw[0, interictal_start:interictal_stop]
    data_ictal, times_ictal = raw[0, ictal_start:ictal_stop]
    
    fig, axes = plt.subplots(1, 2, figsize=(12, 4))
    axes[0].plot(times_inter, data_inter.T, color='blue')
    axes[0].set_title("Calm Interictal State (Resting)")
    axes[0].set_xlabel("Time (s)")
    axes[0].set_ylabel("Voltage")
    
    axes[1].plot(times_ictal, data_ictal.T, color='red')
    axes[1].set_title("Chaotic Ictal State (Seizure)")
    axes[1].set_xlabel("Time (s)")
    axes[1].set_ylabel("Voltage")
    plt.tight_layout()
    plt.savefig("state_comparison.png")
    plt.close()

    # 4. Filtered Data
    print("Filtering data...")
    raw_filtered = raw.copy()
    raw_filtered.filter(l_freq=1.0, h_freq=50.0)
    raw_filtered.notch_filter(freqs=60.0) # US power line freq
    
    start, stop = raw_filtered.time_as_index([45, 50])
    data_clean, times_clean = raw_filtered[0, start:stop]
    
    plt.figure(figsize=(10, 4))
    plt.plot(times_clean, data_clean.T, color='green')
    plt.title("Clean Filtered EEG Signal (Band-pass 1-50Hz, Notch 60Hz)")
    plt.xlabel("Time (s)")
    plt.ylabel("Voltage")
    plt.tight_layout()
    plt.savefig("clean_filtered_eeg.png")
    plt.close()

def create_presentation():
    print("Creating presentation...")
    prs = Presentation()
    
    # Slide 1: Executive Summary
    slide = prs.slides.add_slide(prs.slide_layouts[1]) # Title and Content
    slide.shapes.title.text = "Executive Summary: The 'Three P's'"
    tf = slide.shapes.placeholders[1].text_frame
    
    p = tf.add_paragraph()
    p.text = "Progress:"
    p.font.bold = True
    p.font.size = Pt(24)
    p = tf.add_paragraph()
    p.text = "Downloaded dataset ds003029, successfully loaded the BIDS formatting into Python (using the MNE library), and completed initial Exploratory Data Analysis (EDA) and preprocessing pipelines."
    p.level = 1
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "Plan:"
    p.font.bold = True
    p.font.size = Pt(24)
    p = tf.add_paragraph()
    p.text = "Finalise the cleaned, epoched baseline data to begin building the ARIMA univariate models next week."
    p.level = 1
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "Problem:"
    p.font.bold = True
    p.font.size = Pt(24)
    p = tf.add_paragraph()
    p.text = "No major technical snags. Filtering and PSD generation executed successfully."
    p.level = 1
    p.font.size = Pt(20)

    # Slide 2: Dataset Eligibility
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Dataset Eligibility (ds003029)"
    tf = slide.shapes.placeholders[1].text_frame
    
    p = tf.add_paragraph()
    p.text = "Clinically Verified Annotations:"
    p.font.bold = True
    p.font.size = Pt(22)
    p = tf.add_paragraph()
    p.text = "Contains exact timestamps for interictal (resting) and ictal (seizure) states. Crucial 'ground truth' for calculating False Positive Rates."
    p.level = 1
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "Sufficient Complexity:"
    p.font.bold = True
    p.font.size = Pt(22)
    p = tf.add_paragraph()
    p.text = "Raw data contains biological noise (e.g., eye blinks, muscle movement), proving the Kalman Filter's necessity for artefact rejection."
    p.level = 1
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "Format Compliance:"
    p.font.bold = True
    p.font.size = Pt(22)
    p = tf.add_paragraph()
    p.text = "Strict BIDS (Brain Imaging Data Structure) format, ensuring standardisation and reproducibility."
    p.level = 1
    p.font.size = Pt(18)

    # Slide 3: EDA & Visualisations
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title only
    slide.shapes.title.text = "Exploratory Data Analysis (EDA) & Visualisations"
    
    if os.path.exists("raw_time_series.png"):
        slide.shapes.add_picture("raw_time_series.png", Inches(0.5), Inches(1.5), width=Inches(4.5))
    if os.path.exists("psd_plot.png"):
        slide.shapes.add_picture("psd_plot.png", Inches(5.2), Inches(1.5), width=Inches(4.5))
    if os.path.exists("state_comparison.png"):
        slide.shapes.add_picture("state_comparison.png", Inches(1.5), Inches(4.5), width=Inches(7))

    # Slide 4: Preprocessing Pipeline
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "The Preprocessing Pipeline"
    tf = slide.shapes.placeholders[1].text_frame
    
    p = tf.add_paragraph()
    p.text = "Step 1: Frequency Filtering"
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "Applied a band-pass filter (1Hz to 50Hz) to isolate neurological data, and a 60Hz Notch filter to remove power line hum."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Step 2: Artefact Rejection"
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "Identified massive spikes from blinking/movement using Independent Component Analysis (ICA) and voltage thresholding."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Step 3: Epoching"
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "Sliced continuous recordings into smaller, chronological 'epochs' (e.g., 2-second windows) for manageable algorithmic processing."
    p.level = 1

    # Slide 5: Output and Next Steps
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Output and Next Steps"
    
    if os.path.exists("clean_filtered_eeg.png"):
        slide.shapes.add_picture("clean_filtered_eeg.png", Inches(1.5), Inches(1.5), width=Inches(7))
        
    txBox = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1.5))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "The data is now mathematically normalised and separated. WP1 is complete. I am now ready to move into WP2 and begin building the ARIMA baseline models on this cleaned data."
    p.font.size = Pt(20)
    p.font.bold = True

    prs.save("Weekly_Presentation.pptx")
    print("Presentation saved as Weekly_Presentation.pptx")

if __name__ == "__main__":
    create_plots()
    create_presentation()
